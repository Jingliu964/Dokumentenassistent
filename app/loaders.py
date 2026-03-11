from pathlib import Path
from typing import Iterable, Tuple

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx"}


def iter_documents(paths: Iterable[str], base_dir: str) -> Iterable[Tuple[str, str]]:
    base = Path(base_dir)
    for root in paths:
        root_path = Path(root)
        if not root_path.exists():
            continue
        for path in root_path.rglob("*"):
            if not path.is_file():
                continue
            if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
                continue
            try:
                content = load_text(path)
            except Exception:
                continue
            if not content.strip():
                continue
            try:
                source = str(path.relative_to(base))
            except ValueError:
                source = str(path)
            yield source, content


def load_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md"}:
        return _load_text_file(path)
    if suffix == ".pdf":
        return _load_pdf(path)
    if suffix == ".docx":
        return _load_docx(path)
    if suffix == ".pptx":
        return _load_pptx(path)
    return ""


def _load_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1")


def _load_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    texts = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        texts.append(page_text)
    return "\n".join(texts)


def _load_docx(path: Path) -> str:
    doc = Document(str(path))
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)


def _load_pptx(path: Path) -> str:
    deck = Presentation(str(path))
    texts = []
    for slide in deck.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if text:
                texts.append(text)
    return "\n".join(texts)
